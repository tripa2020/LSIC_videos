"""PPTX text + speaker notes + per-slide PNG render.

Render path: PPTX → PDF (libreoffice headless) → per-page PNG (PyMuPDF/fitz).
Idempotent via slide_index.json cache.
"""

from __future__ import annotations

import shutil
import subprocess
from pathlib import Path

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.slide import Slide as PptxSlide

from src.contracts import DeckIndex, Slide


def extract(path: Path, out_dir: Path, dpi: int = 120) -> DeckIndex:
    out_dir.mkdir(parents=True, exist_ok=True)
    cache = out_dir / "slide_index.json"
    if cache.exists():
        return DeckIndex.model_validate_json(cache.read_text())

    prs = Presentation(str(path))
    slide_text_notes = [
        (_text_from_slide(s), _notes_from_slide(s)) for s in prs.slides
    ]

    pdf_path = out_dir / "deck.pdf"
    if not pdf_path.exists():
        _pptx_to_pdf(path, out_dir, pdf_path)

    doc = fitz.open(str(pdf_path))
    try:
        # Hidden slides differ between python-pptx (counts them) and
        # libreoffice (exports only visible). Use the shorter count; log if mismatched.
        n_pages = doc.page_count
        n_text = len(slide_text_notes)
        if n_pages != n_text:
            print(
                f"  [pptx_handler] {path.name}: pptx={n_text} slides, "
                f"libreoffice={n_pages} pages — using min({n_text}, {n_pages})"
            )
        n = min(n_pages, n_text)
        slides: list[Slide] = []
        for i in range(n):
            txt, notes = slide_text_notes[i]
            page = doc[i]
            png_path = out_dir / f"slide_{i+1:03d}.png"
            page.get_pixmap(dpi=dpi).save(str(png_path))
            slides.append(
                Slide(n=i + 1, text=txt, speaker_notes=notes, png_path=png_path)
            )
    finally:
        doc.close()

    title = (prs.core_properties.title or path.stem).strip() or path.stem
    idx = DeckIndex(asset_id=_asset_id_from(path), title=title, slides=slides)
    cache.write_text(idx.model_dump_json(indent=2))
    return idx


def _text_from_slide(slide: PptxSlide) -> str:
    parts: list[str] = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = "".join(run.text for run in para.runs).strip()
                if t:
                    parts.append(t)
    return "\n".join(parts)


def _notes_from_slide(slide: PptxSlide) -> str:
    if not slide.has_notes_slide:
        return ""
    return (slide.notes_slide.notes_text_frame.text or "").strip()


def _pptx_to_pdf(src: Path, out_dir: Path, target: Path) -> None:
    subprocess.run(
        ["soffice", "--headless", "--convert-to", "pdf",
         "--outdir", str(out_dir), str(src)],
        check=True, capture_output=True,
    )
    produced = out_dir / f"{src.stem}.pdf"
    if produced.exists() and produced != target:
        shutil.move(str(produced), str(target))


def _asset_id_from(path: Path) -> str:
    return path.stem.split("-", 1)[0]
